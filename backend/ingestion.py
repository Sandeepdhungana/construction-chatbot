"""File ingestion and document preparation pipeline with support for all file types."""

from __future__ import annotations

import os
import shutil
from datetime import datetime
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Tuple, Any
import json

import pandas as pd
from fastapi import UploadFile
from langchain_community.document_loaders import PyPDFLoader

try:
    from langchain_community.document_loaders import (
        UnstructuredWordDocumentLoader,
        UnstructuredPowerPointLoader,
    )
    UNSTRUCTURED_AVAILABLE = True
except ImportError:
    UNSTRUCTURED_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# LangChain v1 imports
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

from .tools import data_registry, get_data_registry
from .vectorstore import VectorStoreManager

UPLOAD_ROOT = Path("data/uploads")
UPLOAD_ROOT.mkdir(parents=True, exist_ok=True)

# File registry to track uploaded files and their metadata (per user)
def get_file_registry_path(user_id: int) -> Path:
    """Get file registry path for a specific user."""
    return UPLOAD_ROOT / f"file_registry_user_{user_id}.json"


def _load_file_registry(user_id: int) -> Dict[str, Any]:
    """Load the file registry from disk for a specific user."""
    registry_path = get_file_registry_path(user_id)
    if registry_path.exists():
        with open(registry_path, "r") as f:
            return json.load(f)
    return {}


def _save_file_registry(registry: Dict[str, Any], user_id: int) -> None:
    """Save the file registry to disk for a specific user."""
    registry_path = get_file_registry_path(user_id)
    with open(registry_path, "w") as f:
        json.dump(registry, f, indent=2, default=str)


def _get_file_metadata(file_path: Path) -> Dict[str, Any]:
    """Extract metadata about a file."""
    stat = file_path.stat()
    return {
        "filename": file_path.name,
        "path": str(file_path),
        "size_bytes": stat.st_size,
        "uploaded_at": datetime.utcnow().isoformat(),
    }


def _timestamp_folder(user_id: int, user_email: Optional[str] = None) -> Path:
    """Get timestamped folder for user uploads using email if available."""
    # Get user email if not provided
    if user_email is None:
        from .auth import get_user_by_id
        user = get_user_by_id(user_id)
        if user:
            user_email = user.get("email", "").replace("@", "_at_").replace(".", "_")
        else:
            user_email = f"user_{user_id}"
    
    # Sanitize email for filesystem (replace @ with _at_ and . with _)
    if "@" in user_email:
        user_email = user_email.replace("@", "_at_").replace(".", "_")
    
    user_upload_dir = UPLOAD_ROOT / f"user_{user_email}"
    user_upload_dir.mkdir(parents=True, exist_ok=True)
    return user_upload_dir / datetime.utcnow().strftime("%Y%m%d_%H%M%S")


async def save_uploads(
    files: Iterable[UploadFile], target_dir: Path
) -> List[Path]:
    saved_paths: List[Path] = []
    for file in files:
        destination = target_dir / file.filename
        with destination.open("wb") as f:
            content = await file.read()
            f.write(content)
        saved_paths.append(destination)
    return saved_paths


def _chunk_pdfs(pdf_paths: Iterable[Path], user_id: int) -> List[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1200,
        chunk_overlap=150,
        separators=["\n\n", "\n", ".", " ", ""],
    )
    documents: List[Document] = []
    for path in pdf_paths:
        loader = PyPDFLoader(str(path))
        pages = loader.load()
        chunks = splitter.split_documents(pages)
        for chunk in chunks:
            chunk.metadata.update(
                {
                    "source": "pdf",
                    "filename": path.name,
                    "path": str(path),
                    "user_id": user_id,
                }
            )
        documents.extend(chunks)
    return documents


def _chunk_docx(docx_paths: Iterable[Path], user_id: int) -> List[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1200,
        chunk_overlap=150,
        separators=["\n\n", "\n", ".", " ", ""],
    )
    documents: List[Document] = []
    for path in docx_paths:
        try:
            if UNSTRUCTURED_AVAILABLE:
                loader = UnstructuredWordDocumentLoader(str(path))
                pages = loader.load()
            elif DOCX_AVAILABLE:
                # Fallback to python-docx
                doc = docx.Document(str(path))
                text = "\n".join([para.text for para in doc.paragraphs])
                pages = [Document(page_content=text)]
            else:
                print(f"DOCX loader not available. Install python-docx or unstructured.")
                continue
            
            chunks = splitter.split_documents(pages)
            for chunk in chunks:
                chunk.metadata.update(
                    {
                        "source": "docx",
                        "filename": path.name,
                        "path": str(path),
                        "user_id": user_id,
                    }
                )
            documents.extend(chunks)
        except Exception as e:
            print(f"Error loading DOCX {path}: {e}")
    return documents


def _chunk_pptx(pptx_paths: Iterable[Path], user_id: int) -> List[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1200,
        chunk_overlap=150,
        separators=["\n\n", "\n", ".", " ", ""],
    )
    documents: List[Document] = []
    for path in pptx_paths:
        try:
            if UNSTRUCTURED_AVAILABLE:
                loader = UnstructuredPowerPointLoader(str(path))
                pages = loader.load()
            elif PPTX_AVAILABLE:
                # Fallback to python-pptx
                prs = Presentation(str(path))
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)
                text = "\n".join(text_parts)
                pages = [Document(page_content=text)]
            else:
                print(f"PPTX loader not available. Install python-pptx or unstructured.")
                continue
            
            chunks = splitter.split_documents(pages)
            for chunk in chunks:
                chunk.metadata.update(
                    {
                        "source": "pptx",
                        "filename": path.name,
                        "path": str(path),
                        "user_id": user_id,
                    }
                )
            documents.extend(chunks)
        except Exception as e:
            print(f"Error loading PPTX {path}: {e}")
    return documents


def _extract_image_text(image_path: Path) -> str:
    """Extract text from an image using OCR."""
    if not OCR_AVAILABLE:
        return f"[Image file: {image_path.name} - OCR not available. Install pytesseract and Pillow for text extraction.]"
    try:
        image = Image.open(image_path)
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        print(f"Error extracting text from image {image_path}: {e}")
        return f"[Image file: {image_path.name} - OCR extraction failed: {str(e)}]"


def _chunk_images(image_paths: Iterable[Path], user_id: int) -> List[Document]:
    documents: List[Document] = []
    for path in image_paths:
        text = _extract_image_text(path)
        if text.strip():
            documents.append(
                Document(
                    page_content=text,
                    metadata={
                        "source": "image",
                        "filename": path.name,
                        "path": str(path),
                        "user_id": user_id,
                    }
                )
            )
    return documents


def _analyze_dataframe_structure(df: pd.DataFrame) -> Dict[str, Any]:
    """Analyze a DataFrame structure and infer column types and relationships."""
    structure = {
        "columns": [],
        "row_count": len(df),
        "column_count": len(df.columns),
        "dtypes": {},
        "sample_values": {},
        "null_counts": {},
        "numeric_columns": [],
        "text_columns": [],
        "date_columns": [],
        "potential_keys": [],
    }
    
    for col in df.columns:
        dtype = str(df[col].dtype)
        structure["dtypes"][col] = dtype
        structure["null_counts"][col] = int(df[col].isnull().sum())
        
        # Sample values (non-null)
        sample = df[col].dropna().head(3).tolist()
        structure["sample_values"][col] = [str(v) for v in sample]
        
        col_info = {
            "name": col,
            "dtype": dtype,
            "null_count": int(df[col].isnull().sum()),
            "unique_count": int(df[col].nunique()),
        }
        
        # Categorize columns
        if pd.api.types.is_numeric_dtype(df[col]):
            structure["numeric_columns"].append(col)
            col_info["type_category"] = "numeric"
        elif pd.api.types.is_datetime64_any_dtype(df[col]):
            structure["date_columns"].append(col)
            col_info["type_category"] = "date"
        else:
            structure["text_columns"].append(col)
            col_info["type_category"] = "text"
        
        # Detect potential keys (high uniqueness)
        if df[col].nunique() == len(df) and df[col].notna().all():
            structure["potential_keys"].append(col)
            col_info["is_potential_key"] = True
        else:
            col_info["is_potential_key"] = False
        
        structure["columns"].append(col_info)
    
    return structure


def _load_excel_file(path: Path) -> Dict[str, pd.DataFrame]:
    """Load all sheets from an Excel file."""
    try:
        excel_file = pd.ExcelFile(path)
        sheets = {}
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            sheets[sheet_name] = df
        return sheets
    except Exception as e:
        print(f"Error loading Excel file {path}: {e}")
        return {}


def _csv_documents(path: Path, source_label: str) -> Tuple[List[Document], pd.DataFrame, Dict[str, Any]]:
    """Load CSV and return empty documents list (structured data is handled via data_registry, not embeddings)."""
    try:
        df = pd.read_csv(path)
    except Exception as e:
        print(f"Error reading CSV {path}: {e}")
        df = pd.DataFrame()
    
    structure = _analyze_dataframe_structure(df) if not df.empty else {}
    
    # Don't create documents for structured data - it's queried directly via data_registry
    # Only return empty list since CSV data doesn't need embedding
    docs: List[Document] = []
    
    return docs, df, structure


def _excel_documents(path: Path, source_label: str) -> Tuple[List[Document], Dict[str, pd.DataFrame], Dict[str, Dict[str, Any]]]:
    """Load Excel file and return empty documents list (structured data is handled via data_registry, not embeddings)."""
    sheets = _load_excel_file(path)
    structures: Dict[str, Dict[str, Any]] = {}
    
    for sheet_name, df in sheets.items():
        if df.empty:
            continue
        
        structure = _analyze_dataframe_structure(df)
        structures[sheet_name] = structure
    
    # Don't create documents for structured data - it's queried directly via data_registry
    # Only return empty list since Excel data doesn't need embedding
    all_docs: List[Document] = []
    
    return all_docs, sheets, structures


def _register_dataframe(category: str, name: str, df: pd.DataFrame, structure: Optional[Dict[str, Any]] = None, user_id: Optional[int] = None) -> None:
    """Register a DataFrame with optional structure metadata."""
    user_registry = get_data_registry(user_id) if user_id is not None else data_registry
    user_registry.register(category, name, df, structure, user_id=user_id)


def _get_file_type(file_path: Path) -> str:
    """Determine file type from extension."""
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return "pdf"
    elif ext in [".doc", ".docx"]:
        return "docx"
    elif ext in [".ppt", ".pptx"]:
        return "pptx"
    elif ext == ".csv":
        return "csv"
    elif ext in [".xls", ".xlsx", ".xlsm"]:
        return "excel"
    elif ext in [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff"]:
        return "image"
    else:
        return "unknown"


async def ingest_files(
    files: List[UploadFile],
    user_id: int,
    manager: Optional[VectorStoreManager] = None,
) -> dict:
    """Generic file ingestion supporting all file types."""
    # Create user-specific vectorstore manager if not provided
    if manager is None:
        from .vectorstore import get_vector_manager
        manager = get_vector_manager(user_id)
    
    # Get user email for folder naming
    from .auth import get_user_by_id
    user = get_user_by_id(user_id)
    user_email = user.get("email") if user else None
    
    # Get user-specific data registry
    user_registry = get_data_registry(user_id)
    
    ingestion_dir = _timestamp_folder(user_id, user_email=user_email)
    ingestion_dir.mkdir(parents=True, exist_ok=True)
    
    saved_paths = await save_uploads(files, ingestion_dir)
    
    all_docs: List[Document] = []
    file_registry = _load_file_registry(user_id)
    
    # Group files by type
    pdfs = []
    docxs = []
    pptxs = []
    csvs = []
    excels = []
    images = []
    
    for path in saved_paths:
        file_type = _get_file_type(path)
        file_meta = _get_file_metadata(path)
        file_id = str(path.relative_to(UPLOAD_ROOT))
        file_registry[file_id] = {
            **file_meta,
            "type": file_type,
            "user_id": user_id,
        }
        
        if file_type == "pdf":
            pdfs.append(path)
        elif file_type == "docx":
            docxs.append(path)
        elif file_type == "pptx":
            pptxs.append(path)
        elif file_type == "csv":
            csvs.append(path)
        elif file_type == "excel":
            excels.append(path)
        elif file_type == "image":
            images.append(path)
    
    # Process PDFs
    if pdfs:
        all_docs.extend(_chunk_pdfs(pdfs, user_id))
    
    # Process DOCX files
    if docxs:
        all_docs.extend(_chunk_docx(docxs, user_id))
    
    # Process PPTX files
    if pptxs:
        all_docs.extend(_chunk_pptx(pptxs, user_id))
    
    # Process images
    if images:
        all_docs.extend(_chunk_images(images, user_id))
    
    # Process CSV files (generic, no assumptions)
    csv_count = 0
    import logging
    ingestion_logger = logging.getLogger(__name__)
    
    for path in csvs:
        ingestion_logger.info(f"📄 Processing CSV: {path.name}")
        docs, df, structure = _csv_documents(path, "csv")
        all_docs.extend(docs)
        table_name = f"{path.stem}"
        ingestion_logger.info(f"📊 Registering table '{table_name}' with {len(df)} rows, {len(df.columns)} columns")
        _register_dataframe("spreadsheet", table_name, df, structure, user_id=user_id)
        csv_count += 1
        
        # Verify registration
        registered_df = user_registry.get("spreadsheet", table_name, user_id=user_id)
        if registered_df is not None:
            ingestion_logger.info(f"✅ Successfully registered table '{table_name}'")
        else:
            ingestion_logger.error(f"❌ Failed to register table '{table_name}'")
        
        # Update registry with structure info
        file_id = str(path.relative_to(UPLOAD_ROOT))
        file_registry[file_id]["structure"] = structure
        file_registry[file_id]["table_name"] = table_name
    
    # Process Excel files (generic, all sheets)
    excel_count = 0
    for path in excels:
        ingestion_logger.info(f"📄 Processing Excel: {path.name}")
        docs, sheets, structures = _excel_documents(path, "excel")
        all_docs.extend(docs)
        
        for sheet_name, df in sheets.items():
            table_name = f"{path.stem}_{sheet_name}"
            structure = structures.get(sheet_name, {})
            ingestion_logger.info(f"📊 Registering table '{table_name}' (sheet: {sheet_name}) with {len(df)} rows, {len(df.columns)} columns")
            _register_dataframe("spreadsheet", table_name, df, structure, user_id=user_id)
            excel_count += 1
            
            # Verify registration
            registered_df = user_registry.get("spreadsheet", table_name, user_id=user_id)
            if registered_df is not None:
                ingestion_logger.info(f"✅ Successfully registered table '{table_name}'")
            else:
                ingestion_logger.error(f"❌ Failed to register table '{table_name}'")
        
        # Update registry
        file_id = str(path.relative_to(UPLOAD_ROOT))
        file_registry[file_id]["sheets"] = list(sheets.keys())
        file_registry[file_id]["structures"] = structures
    
    # Save file registry
    _save_file_registry(file_registry, user_id)
    
    # Add all documents to vectorstore
    added = manager.add_documents(all_docs) if all_docs else 0
    
    # Log final registry state
    user_registry = get_data_registry(user_id)
    final_summary = user_registry.summary(user_id=user_id)
    ingestion_logger.info(f"📊 Final registry summary: {final_summary}")
    ingestion_logger.info(f"📊 Total tables registered: {sum(len(tables) for tables in final_summary.values())}")
    
    return {
        "documents_added": added,
        "files_uploaded": len(saved_paths),
        "pdf_count": len(pdfs),
        "docx_count": len(docxs),
        "pptx_count": len(pptxs),
        "csv_count": csv_count,
        "excel_count": excel_count,
        "image_count": len(images),
        "dataframes": final_summary,
        "storage_path": str(ingestion_dir),
    }


async def ingest_payload(
    pdf_files: Iterable[UploadFile],
    materials_files: Iterable[UploadFile],
    workforce_files: Iterable[UploadFile],
    user_id: int,
    manager: Optional[VectorStoreManager] = None,
) -> dict:
    """Legacy function for backward compatibility."""
    all_files = list(pdf_files) + list(materials_files) + list(workforce_files)
    return await ingest_files(all_files, user_id, manager)


def reload_tables_from_registry(user_id: int, data_registry_instance: Optional[Any] = None) -> int:
    """Reload all spreadsheet tables from the file registry for a specific user."""
    import logging
    ingestion_logger = logging.getLogger(__name__)
    
    registry = _load_file_registry(user_id)
    reloaded_count = 0
    
    ingestion_logger.info(f"🔄 Reloading tables from registry for user {user_id} ({len(registry)} files)...")
    
    # Use provided registry instance or get a new one
    if data_registry_instance is None:
        from .tools import get_data_registry
        data_registry_instance = get_data_registry(user_id, reload=False)
    
    # Get user email to find the correct user folder
    from .auth import get_user_by_id
    user = get_user_by_id(user_id)
    user_email = None
    if user:
        user_email = user.get("email", "").replace("@", "_at_").replace(".", "_")
    
    for file_id, metadata in registry.items():
        # CRITICAL: Only reload files for this user - strict check
        if metadata.get("user_id") != user_id:
            ingestion_logger.debug(f"⏭️  Skipping file {file_id} - belongs to user {metadata.get('user_id')}, not {user_id}")
            continue
            
        file_path = UPLOAD_ROOT / file_id
        
        # If file doesn't exist at the expected path, try to find it ONLY in current user's folder
        if not file_path.exists():
            file_name = Path(file_id).name
            
            # Only search in the current user's folder (by email or user_id)
            user_folders_to_check = []
            if user_email:
                user_folders_to_check.append(f"user_{user_email}")
            # Also check old format for backward compatibility
            user_folders_to_check.append(f"user_{user_id}")
            
            found = False
            for folder_name in user_folders_to_check:
                user_folder = UPLOAD_ROOT / folder_name
                if user_folder.exists() and user_folder.is_dir():
                    # Search in timestamp subdirectories
                    for timestamp_dir in user_folder.iterdir():
                        if timestamp_dir.is_dir():
                            potential_file = timestamp_dir / file_name
                            if potential_file.exists():
                                file_path = potential_file
                                found = True
                                ingestion_logger.info(f"📁 Found file in {folder_name}/{timestamp_dir.name}/")
                                break
                        if found:
                            break
                if found:
                    break
            
            # If still not found, skip this file - don't search other users' folders
            if not file_path.exists():
                ingestion_logger.warning(f"⚠️  File not found for user {user_id}: {file_id}, skipping...")
                continue
        
        file_type = metadata.get("type")
        if file_type == "csv":
            try:
                ingestion_logger.info(f"📄 Reloading CSV: {file_path.name}")
                _, df, structure = _csv_documents(file_path, "csv")  # Ignore docs, we don't embed CSV
                table_name = metadata.get("table_name") or file_path.stem
                data_registry_instance.register("spreadsheet", table_name, df, structure, user_id=user_id)
                reloaded_count += 1
                ingestion_logger.info(f"✅ Reloaded table '{table_name}'")
            except Exception as e:
                ingestion_logger.error(f"❌ Failed to reload CSV {file_path.name}: {e}")
        
        elif file_type == "excel":
            try:
                ingestion_logger.info(f"📄 Reloading Excel: {file_path.name}")
                _, sheets, structures = _excel_documents(file_path, "excel")  # Ignore docs, we don't embed Excel
                for sheet_name, df in sheets.items():
                    table_name = f"{file_path.stem}_{sheet_name}"
                    structure = structures.get(sheet_name, {})
                    data_registry_instance.register("spreadsheet", table_name, df, structure, user_id=user_id)
                    reloaded_count += 1
                    ingestion_logger.info(f"✅ Reloaded table '{table_name}'")
            except Exception as e:
                ingestion_logger.error(f"❌ Failed to reload Excel {file_path.name}: {e}")
    
    ingestion_logger.info(f"✅ Reloaded {reloaded_count} table(s) from registry")
    return reloaded_count


def get_uploaded_files(user_id: int) -> List[Dict[str, Any]]:
    """Get list of all uploaded files with metadata for a specific user."""
    registry = _load_file_registry(user_id)
    files = []
    for file_id, metadata in registry.items():
        file_path = UPLOAD_ROOT / file_id
        if file_path.exists() and metadata.get("user_id") == user_id:
            files.append({
                "id": file_id,
                **metadata,
            })
    return sorted(files, key=lambda x: x.get("uploaded_at", ""), reverse=True)


def delete_file(file_id: str, user_id: int, manager: Optional[VectorStoreManager] = None) -> bool:
    """Delete a file and remove it from registry, data registry, and vectorstore."""
    import logging
    ingestion_logger = logging.getLogger(__name__)
    
    # Create user-specific vectorstore manager if not provided
    if manager is None:
        from .vectorstore import get_vector_manager
        manager = get_vector_manager(user_id)
    
    registry = _load_file_registry(user_id)
    if file_id not in registry:
        return False
    
    # Verify the file belongs to this user
    if registry[file_id].get("user_id") != user_id:
        return False
    
    file_path = UPLOAD_ROOT / file_id
    metadata = registry[file_id]
    filename = metadata.get("filename", file_path.name)
    file_path_str = str(file_path) if file_path.exists() else None
    
    ingestion_logger.info(f"🗑️  Deleting file: {filename} (ID: {file_id})")
    
    # Delete from disk
    if file_path.exists():
        file_path.unlink()
        ingestion_logger.info(f"✅ Deleted file from disk: {file_path}")
    
    # Remove from data registry if it's a spreadsheet
    if metadata.get("type") in ["csv", "excel"]:
        user_registry = get_data_registry(user_id)
        table_name = metadata.get("table_name")
        if table_name:
            user_registry.unregister("spreadsheet", table_name, user_id=user_id)
            ingestion_logger.info(f"✅ Removed table '{table_name}' from data registry")
        elif metadata.get("type") == "excel":
            # Remove all sheets
            for sheet_name in metadata.get("sheets", []):
                table_name = f"{Path(file_id).stem}_{sheet_name}"
                user_registry.unregister("spreadsheet", table_name, user_id=user_id)
                ingestion_logger.info(f"✅ Removed table '{table_name}' from data registry")
    
    # Delete from vectorstore by filename and path
    deleted_count = 0
    if file_path_str:
        # Try deleting by path first (most specific)
        deleted_count = manager.delete_by_metadata(path=file_path_str)
        if deleted_count == 0:
            # Fallback to filename if path doesn't match
            deleted_count = manager.delete_by_metadata(filename=filename)
    
    if deleted_count > 0:
        ingestion_logger.info(f"✅ Deleted {deleted_count} document(s) from vectorstore")
    else:
        ingestion_logger.warning(f"⚠️  No documents found in vectorstore for file: {filename}")
    
    # Remove from file registry
    del registry[file_id]
    _save_file_registry(registry, user_id)
    ingestion_logger.info(f"✅ Removed file from registry")
    
    return True


def reset_uploads(user_id: Optional[int] = None) -> None:
    """Remove all uploaded files and reset registry/vectorstore.
    
    Args:
        user_id: If provided, only resets data for that user. Otherwise resets all.
    """
    if UPLOAD_ROOT.exists():
        shutil.rmtree(UPLOAD_ROOT)
    UPLOAD_ROOT.mkdir(parents=True, exist_ok=True)
    if user_id is not None:
        data_registry.clear_all(user_id=user_id)
        # Note: VectorStoreManager is per-user, so clearing is handled per user
        # This function is mainly for testing/cleanup
    else:
        data_registry.clear_all()
